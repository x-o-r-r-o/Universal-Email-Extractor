import os
import re
import argparse
import logging
import zipfile
import tarfile
import sqlite3
import pdfplumber
import pytesseract
from PIL import Image
import docx
import openpyxl
import xlrd
import pandas as pd
import filetype
import chardet
from striprtf.striprtf import rtf_to_text
from odf import text, teletype
from odf.opendocument import load as odf_load
import extract_msg
from tqdm import tqdm
import string

def setup_logger(logfile):
    logging.basicConfig(
        filename=logfile,
        level=logging.INFO,
        format='%(asctime)s [%(levelname)s] %(message)s'
    )

def extract_emails_from_text(text):
    # Stricter regex for emails with TLD at least 2 chars, no trailing image/file extensions
    email_regex = r'\b[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z]{2,}\b'
    all_matches = set(re.findall(email_regex, text or ""))

    # Remove emails that end with image/file extensions or numbers after @
    file_exts = [
        '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.svg', '.webp', '.ico',
        '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.zip', '.rar'
    ]
    filtered = set()
    for email in all_matches:
        # Remove trailing punctuation
        email_clean = email.strip().strip(string.punctuation)
        # Exclude if ends with known file extension or has numbers after @
        domain_part = email_clean.split('@')[-1]
        if any(domain_part.lower().endswith(ext) for ext in file_exts):
            continue
        if re.search(r'@\d', email_clean):
            continue
        filtered.add(email_clean)
    return filtered

def read_text_file(file_path):
    with open(file_path, 'rb') as f:
        raw = f.read()
    enc = chardet.detect(raw)['encoding'] or 'utf-8'
    try:
        text = raw.decode(enc, errors='ignore')
    except Exception:
        text = raw.decode('utf-8', errors='ignore')
    return extract_emails_from_text(text)

def read_csv_file(file_path):
    try:
        try:
            df = pd.read_csv(file_path, dtype=str, encoding='utf-8')
        except UnicodeDecodeError:
            with open(file_path, 'rb') as f:
                raw = f.read()
            enc = chardet.detect(raw)['encoding'] or 'utf-8'
            df = pd.read_csv(file_path, dtype=str, encoding=enc)
        return extract_emails_from_text(df.to_string())
    except Exception as e:
        logging.error(f"CSV processing failed for {file_path}: {e}")
        return set()

def read_xls_file(file_path):
    try:
        wb = xlrd.open_workbook(file_path)
        emails = set()
        for sheet in wb.sheets():
            for row in range(sheet.nrows):
                for val in sheet.row_values(row):
                    if isinstance(val, str):
                        emails.update(extract_emails_from_text(val))
        return emails
    except Exception as e:
        logging.error(f"XLS processing failed for {file_path}: {e}")
        return set()

def read_xlsx_file(file_path):
    emails = set()
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True)
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if isinstance(cell, str):
                        emails.update(extract_emails_from_text(cell))
    except Exception as e:
        logging.error(f"XLSX processing failed for {file_path}: {e}")
    return emails

def read_docx_file(file_path):
    try:
        doc = docx.Document(file_path)
        text = '\n'.join([para.text for para in doc.paragraphs])
        return extract_emails_from_text(text)
    except Exception as e:
        logging.error(f"DOCX processing failed for {file_path}: {e}")
        return set()

def read_doc_file(file_path):
    try:
        import textract
        text = textract.process(file_path).decode('utf-8', errors='ignore')
        return extract_emails_from_text(text)
    except Exception as e:
        logging.error(f"DOC processing failed for {file_path}: {e}")
        return set()

def read_rtf_file(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            rtf = f.read()
        text = rtf_to_text(rtf)
        return extract_emails_from_text(text)
    except Exception as e:
        logging.error(f"RTF processing failed for {file_path}: {e}")
        return set()

def read_odt_file(file_path):
    try:
        odt = odf_load(file_path)
        texts = odt.getElementsByType(text.P)
        all_text = '\n'.join([teletype.extractText(t) for t in texts])
        return extract_emails_from_text(all_text)
    except Exception as e:
        logging.error(f"ODT processing failed for {file_path}: {e}")
        return set()

def read_pdf_file(file_path):
    emails = set()
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    emails.update(extract_emails_from_text(text))
                else:
                    img = page.to_image(resolution=300)
                    text = pytesseract.image_to_string(img.original)
                    emails.update(extract_emails_from_text(text))
    except Exception as e:
        logging.error(f"PDF processing failed for {file_path}: {e}")
    return emails

def read_image_file(file_path):
    try:
        text = pytesseract.image_to_string(Image.open(file_path))
        return extract_emails_from_text(text)
    except Exception as e:
        logging.error(f"OCR failed for {file_path}: {e}")
        return set()

def read_sqlite_file(file_path):
    emails = set()
    try:
        conn = sqlite3.connect(file_path)
        cursor = conn.cursor()
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = cursor.fetchall()
        for table in tables:
            table_name = table[0]
            try:
                cursor.execute(f"PRAGMA table_info({table_name})")
                columns = [col[1] for col in cursor.fetchall()]
                for col in columns:
                    try:
                        cursor.execute(f"SELECT {col} FROM {table_name}")
                        for row in cursor.fetchall():
                            val = row[0]
                            if isinstance(val, str):
                                emails.update(extract_emails_from_text(val))
                    except Exception:
                        continue
            except Exception:
                continue
        conn.close()
    except Exception as e:
        logging.error(f"SQLite processing failed for {file_path}: {e}")
    return emails

def read_sql_file(file_path):
    return read_text_file(file_path)

def read_pptx_file(file_path):
    try:
        from pptx import Presentation
        emails = set()
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    emails.update(extract_emails_from_text(shape.text))
        return emails
    except Exception as e:
        logging.error(f"PPTX processing failed for {file_path}: {e}")
        return set()

def read_ppt_file(file_path):
    try:
        import textract
        text = textract.process(file_path).decode('utf-8', errors='ignore')
        return extract_emails_from_text(text)
    except Exception as e:
        logging.error(f"PPT processing failed for {file_path}: {e}")
        return set()

def read_msg_file(file_path):
    try:
        msg = extract_msg.Message(file_path)
        body = msg.body or ""
        subj = msg.subject or ""
        attachments = ""
        for att in msg.attachments:
            attachments += att.longFilename + " "
        return extract_emails_from_text(body + subj + attachments)
    except Exception as e:
        logging.error(f"MSG processing failed for {file_path}: {e}")
        return set()

def read_eml_file(file_path):
    try:
        import email
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            msg = email.message_from_file(f)
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == 'text/plain':
                    body += part.get_payload(decode=True).decode('utf-8', errors='ignore')
        else:
            body = msg.get_payload(decode=True).decode('utf-8', errors='ignore')
        return extract_emails_from_text(body)
    except Exception as e:
        logging.error(f"EML processing failed for {file_path}: {e}")
        return set()

def read_mdb_file(file_path):
    logging.warning(f"MDB/ACCDB handler stub for {file_path}. Install pyodbc/msaccessdb for full support.")
    return set()

def read_archive(file_path, temp_dir):
    emails = set()
    try:
        if zipfile.is_zipfile(file_path):
            with zipfile.ZipFile(file_path, 'r') as archive:
                archive.extractall(temp_dir)
        elif tarfile.is_tarfile(file_path):
            with tarfile.open(file_path, 'r:*') as archive:
                archive.extractall(temp_dir)
        else:
            return emails
        for root, _, files in os.walk(temp_dir):
            for file in files:
                full_path = os.path.join(root, file)
                emails.update(process_file(full_path, temp_dir))
    except Exception as e:
        logging.error(f"Archive extraction failed for {file_path}: {e}")
    return emails

def process_file(file_path, temp_dir):
    ext = os.path.splitext(file_path)[1].lower()
    handlers = {
        # Text/Config/Code
        '.txt': read_text_file, '.log': read_text_file, '.ini': read_text_file, '.inf': read_text_file,
        '.html': read_text_file, '.htm': read_text_file, '.asp': read_text_file, '.aspx': read_text_file,
        '.php': read_text_file, '.js': read_text_file, '.json': read_text_file, '.xml': read_text_file,
        '.yaml': read_text_file, '.yml': read_text_file, '.md': read_text_file,
        # Spreadsheets
        '.csv': read_csv_file, '.xls': read_xls_file, '.xlsx': read_xlsx_file, '.xlsm': read_xlsx_file,
        '.ods': read_odt_file,
        # Office Docs
        '.docx': read_docx_file, '.docm': read_docx_file, '.doc': read_doc_file, '.rtf': read_rtf_file,
        '.odt': read_odt_file,
        # Presentations
        '.pptx': read_pptx_file, '.ppt': read_ppt_file,
        # PDFs
        '.pdf': read_pdf_file,
        # Images
        '.jpg': read_image_file, '.jpeg': read_image_file, '.png': read_image_file, '.bmp': read_image_file,
        '.tiff': read_image_file, '.gif': read_image_file,
        # Databases
        '.sqlite': read_sqlite_file, '.sqlite3': read_sqlite_file, '.db': read_sqlite_file, '.sql': read_sql_file,
        '.mdb': read_mdb_file, '.accdb': read_mdb_file,
        # Emails
        '.eml': read_eml_file, '.msg': read_msg_file,
        # Archives
        '.zip': lambda f: read_archive(f, temp_dir), '.tar': lambda f: read_archive(f, temp_dir),
        '.gz': lambda f: read_archive(f, temp_dir), '.rar': lambda f: read_archive(f, temp_dir),
    }
    if ext in handlers:
        return handlers[ext](file_path)
    kind = filetype.guess(file_path)
    if kind and kind.mime.startswith('image'):
        return read_image_file(file_path)
    try:
        return read_text_file(file_path)
    except Exception:
        return set()

def is_compatible_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    compatible_exts = [
        '.txt', '.log', '.ini', '.inf', '.html', '.htm', '.asp', '.aspx', '.php', '.js', '.json', '.xml', '.yaml', '.yml', '.md',
        '.csv', '.xls', '.xlsx', '.xlsm', '.ods',
        '.docx', '.docm', '.doc', '.rtf', '.odt',
        '.pptx', '.ppt',
        '.pdf',
        '.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif',
        '.sqlite', '.sqlite3', '.db', '.sql', '.mdb', '.accdb',
        '.eml', '.msg',
        '.zip', '.tar', '.gz', '.rar'
    ]
    return ext in compatible_exts

def scan_folder(folder, output_file, log_file):
    setup_logger(log_file)
    all_emails = set()
    temp_dir = os.path.join(folder, "_temp_extract")
    os.makedirs(temp_dir, exist_ok=True)
    file_list = []
    for root, _, files in os.walk(folder):
        for file in files:
            if '_temp_extract' in root:
                continue
            path = os.path.join(root, file)
            file_list.append(path)
    total_files = len(file_list)
    compatible_files = [f for f in file_list if is_compatible_file(f)]
    total_compatible = len(compatible_files)

    print(f"Total files found: {total_files}")
    print(f"Compatible files for extraction: {total_compatible}")

    if total_compatible == 0:
        print("No compatible files found for extraction. Exiting.")
        logging.error("No compatible files found for extraction. Exiting.")
        return

    found_files = 0
    with open(output_file, 'a', buffering=1) as f_out, tqdm(total=total_compatible, desc="Extracting emails", ncols=80) as pbar:
        for idx, path in enumerate(compatible_files, 1):
            print(f"Processing file {idx}/{total_compatible}: {path}")
            logging.info(f"Processing file {idx}/{total_compatible}: {path}")
            try:
                emails = process_file(path, temp_dir)
                new_emails = emails - all_emails
                if new_emails:
                    found_files += 1
                    print(f"  Found emails in {path}: {new_emails}")
                    logging.info(f"Found in {path}: {new_emails}")
                    for email in sorted(new_emails):
                        f_out.write(email + '\n')
                        f_out.flush()
                all_emails.update(new_emails)
            except Exception as e:
                print(f"  Error processing {path}: {e}")
                logging.error(f"Failed to process {path}: {e}")
            pbar.update(1)
    logging.info(f"Extraction complete. Unique emails found: {len(all_emails)}")
    import shutil
    shutil.rmtree(temp_dir, ignore_errors=True)
    print(f"Extraction complete. Unique emails found: {len(all_emails)}. See {output_file}.")
    print(f"Files with emails found: {found_files} / {total_compatible}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Universal Production-Ready Email Extractor (Real-Time Output, Strict Filtering)")
    parser.add_argument("folder", help="Folder to scan")
    parser.add_argument("-o", "--output", default="emails_found.txt", help="Output file")
    parser.add_argument("-l", "--log", default="extractor.log", help="Log file")
    args = parser.parse_args()
    scan_folder(args.folder, args.output, args.log)