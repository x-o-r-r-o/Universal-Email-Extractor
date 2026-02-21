import os
import re
import argparse
import logging
import zipfile
import tarfile
import sqlite3
import pdfplumber
import pytesseract
from PIL import Image
import docx
import openpyxl
import xlrd
import pandas as pd
import filetype
import chardet
from striprtf.striprtf import rtf_to_text
from odf import text, teletype
from odf.opendocument import load as odf_load
import extract_msg
from tqdm import tqdm
import string
from urllib.parse import urlparse

FORBIDDEN_WORDS = [
    'user', 'users', 'test', 'example', 'demo', 'sample', 'dummy', 'temp', 'trial', 'no-reply', 'noreply'
]

def setup_logger(logfile):
    logging.basicConfig(
        filename=logfile,
        level=logging.INFO,
        format='%(asctime)s [%(levelname)s] %(message)s'
    )

def load_disposable_domains(filepath='disposable_domains.txt'):
    disposable_domains = set()
    with open(filepath, 'r') as f:
        for line in f:
            domain = line.strip().lower()
            if domain and not domain.startswith('#'):
                disposable_domains.add(domain)
    return disposable_domains

def load_blocked_domains(filepath='blocked_domains.txt'):
    blocked_domains = set()
    with open(filepath, 'r', encoding='utf-8') as f:
        for line in f:
            domain = line.strip().lower()
            if domain and not domain.startswith('#'):
                blocked_domains.add(domain)
    return blocked_domains

def extract_emails_from_text(text, disposable_domains):
    email_regex = r'\b[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z]{2,}\b'
    all_matches = set(re.findall(email_regex, text or ""))

    file_exts = [
        '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.svg', '.webp', '.ico',
        '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.zip', '.rar'
    ]
    filtered = set()
    for email in all_matches:
        email_clean = email.strip().strip(string.punctuation)
        local, _, domain = email_clean.lower().partition('@')
        if any(domain.endswith(ext) for ext in file_exts):
            continue
        if re.search(r'@\d', email_clean):
            continue
        if any(word in local or word in domain for word in FORBIDDEN_WORDS):
            continue
        domain_only = domain.split(':')[0].split('/')[0]
        if domain_only in disposable_domains:
            continue
        filtered.add(email_clean)
    return filtered

def extract_urls_from_text(text, blocked_domains):
    url_regex = r'https?://[^\s"<>\]]+'
    all_urls = set(re.findall(url_regex, text or ""))
    filtered = set()
    for url in all_urls:
        try:
            parsed = urlparse(url)
            domain = parsed.netloc.lower()
            if domain.startswith('www.'):
                domain = domain[4:]
            if domain in blocked_domains:
                continue
            filtered.add(url)
        except Exception:
            continue
    return filtered

def read_text_file(file_path, disposable_domains):
    with open(file_path, 'rb') as f:
        raw = f.read()
    enc = chardet.detect(raw)['encoding'] or 'utf-8'
    try:
        text = raw.decode(enc, errors='ignore')
    except Exception:
        text = raw.decode('utf-8', errors='ignore')
    return extract_emails_from_text(text, disposable_domains), text

def read_csv_file(file_path, disposable_domains):
    try:
        try:
            df = pd.read_csv(file_path, dtype=str, encoding='utf-8')
        except UnicodeDecodeError:
            with open(file_path, 'rb') as f:
                raw = f.read()
            enc = chardet.detect(raw)['encoding'] or 'utf-8'
            df = pd.read_csv(file_path, dtype=str, encoding=enc)
        text = df.to_string()
        return extract_emails_from_text(text, disposable_domains), text
    except Exception as e:
        logging.error(f"CSV processing failed for {file_path}: {e}")
        return set(), ""

def read_xls_file(file_path, disposable_domains):
    try:
        wb = xlrd.open_workbook(file_path)
        emails = set()
        text = ""
        for sheet in wb.sheets():
            for row in range(sheet.nrows):
                for val in sheet.row_values(row):
                    if isinstance(val, str):
                        emails.update(extract_emails_from_text(val, disposable_domains))
                        text += val + "\n"
        return emails, text
    except Exception as e:
        logging.error(f"XLS processing failed for {file_path}: {e}")
        return set(), ""

def read_xlsx_file(file_path, disposable_domains):
    emails = set()
    text = ""
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True)
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if isinstance(cell, str):
                        emails.update(extract_emails_from_text(cell, disposable_domains))
                        text += cell + "\n"
    except Exception as e:
        logging.error(f"XLSX processing failed for {file_path}: {e}")
    return emails, text

def read_docx_file(file_path, disposable_domains):
    try:
        doc = docx.Document(file_path)
        text = '\n'.join([para.text for para in doc.paragraphs])
        return extract_emails_from_text(text, disposable_domains), text
    except Exception as e:
        logging.error(f"DOCX processing failed for {file_path}: {e}")
        return set(), ""

def read_doc_file(file_path, disposable_domains):
    try:
        import textract
        text = textract.process(file_path).decode('utf-8', errors='ignore')
        return extract_emails_from_text(text, disposable_domains), text
    except Exception as e:
        logging.error(f"DOC processing failed for {file_path}: {e}")
        return set(), ""

def read_rtf_file(file_path, disposable_domains):
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            rtf = f.read()
        text = rtf_to_text(rtf)
        return extract_emails_from_text(text, disposable_domains), text
    except Exception as e:
        logging.error(f"RTF processing failed for {file_path}: {e}")
        return set(), ""

def read_odt_file(file_path, disposable_domains):
    try:
        odt = odf_load(file_path)
        texts = odt.getElementsByType(text.P)
        all_text = '\n'.join([teletype.extractText(t) for t in texts])
        return extract_emails_from_text(all_text, disposable_domains), all_text
    except Exception as e:
        logging.error(f"ODT processing failed for {file_path}: {e}")
        return set(), ""

def read_pdf_file(file_path, disposable_domains):
    emails = set()
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    emails.update(extract_emails_from_text(t, disposable_domains))
                    text += t + "\n"
                else:
                    img = page.to_image(resolution=300)
                    t = pytesseract.image_to_string(img.original)
                    emails.update(extract_emails_from_text(t, disposable_domains))
                    text += t + "\n"
    except Exception as e:
        logging.error(f"PDF processing failed for {file_path}: {e}")
    return emails, text

def read_image_file(file_path, disposable_domains):
    try:
        text = pytesseract.image_to_string(Image.open(file_path))
        return extract_emails_from_text(text, disposable_domains), text
    except Exception as e:
        logging.error(f"OCR failed for {file_path}: {e}")
        return set(), ""

def read_sqlite_file(file_path, disposable_domains):
    emails = set()
    text = ""
    try:
        conn = sqlite3.connect(file_path)
        cursor = conn.cursor()
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = cursor.fetchall()
        for table in tables:
            table_name = table[0]
            try:
                cursor.execute(f"PRAGMA table_info({table_name})")
                columns = [col[1] for col in cursor.fetchall()]
                for col in columns:
                    try:
                        cursor.execute(f"SELECT {col} FROM {table_name}")
                        for row in cursor.fetchall():
                            val = row[0]
                            if isinstance(val, str):
                                emails.update(extract_emails_from_text(val, disposable_domains))
                                text += val + "\n"
                    except Exception:
                        continue
            except Exception:
                continue
        conn.close()
    except Exception as e:
        logging.error(f"SQLite processing failed for {file_path}: {e}")
    return emails, text

def read_sql_file(file_path, disposable_domains):
    return read_text_file(file_path, disposable_domains)

def read_pptx_file(file_path, disposable_domains):
    try:
        from pptx import Presentation
        emails = set()
        text = ""
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    emails.update(extract_emails_from_text(shape.text, disposable_domains))
                    text += shape.text + "\n"
        return emails, text
    except Exception as e:
        logging.error(f"PPTX processing failed for {file_path}: {e}")
        return set(), ""

def read_ppt_file(file_path, disposable_domains):
    try:
        import textract
        text = textract.process(file_path).decode('utf-8', errors='ignore')
        return extract_emails_from_text(text, disposable_domains), text
    except Exception as e:
        logging.error(f"PPT processing failed for {file_path}: {e}")
        return set(), ""

def read_msg_file(file_path, disposable_domains):
    try:
        msg = extract_msg.Message(file_path)
        body = msg.body or ""
        subj = msg.subject or ""
        attachments = ""
        for att in msg.attachments:
            attachments += att.longFilename + " "
        text = body + subj + attachments
        return extract_emails_from_text(text, disposable_domains), text
    except Exception as e:
        logging.error(f"MSG processing failed for {file_path}: {e}")
        return set(), ""

def read_eml_file(file_path, disposable_domains):
    try:
        import email
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            msg = email.message_from_file(f)
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == 'text/plain':
                    body += part.get_payload(decode=True).decode('utf-8', errors='ignore')
        else:
            body = msg.get_payload(decode=True).decode('utf-8', errors='ignore')
        return extract_emails_from_text(body, disposable_domains), body
    except Exception as e:
        logging.error(f"EML processing failed for {file_path}: {e}")
        return set(), ""

def read_mdb_file(file_path, disposable_domains):
    logging.warning(f"MDB/ACCDB handler stub for {file_path}. Install pyodbc/msaccessdb for full support.")
    return set(), ""

def read_archive(file_path, temp_dir, disposable_domains):
    emails = set()
    text = ""
    try:
        if zipfile.is_zipfile(file_path):
            with zipfile.ZipFile(file_path, 'r') as archive:
                archive.extractall(temp_dir)
        elif tarfile.is_tarfile(file_path):
            with tarfile.open(file_path, 'r:*') as archive:
                archive.extractall(temp_dir)
        else:
            return emails, text
        for root, _, files in os.walk(temp_dir):
            for file in files:
                full_path = os.path.join(root, file)
                e, t = process_file(full_path, temp_dir, disposable_domains)
                emails.update(e)
                text += t + "\n"
    except Exception as e:
        logging.error(f"Archive extraction failed for {file_path}: {e}")
    return emails, text

def is_compatible_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    compatible_exts = [
        '.txt', '.log', '.ini', '.inf', '.html', '.htm', '.asp', '.aspx', '.php', '.js', '.json', '.xml', '.yaml', '.yml', '.md',
        '.csv', '.xls', '.xlsx', '.xlsm', '.ods',
        '.docx', '.docm', '.doc', '.rtf', '.odt',
        '.pptx', '.ppt',
        '.pdf',
        '.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif',
        '.sqlite', '.sqlite3', '.db', '.sql', '.mdb', '.accdb',
        '.eml', '.msg',
        '.zip', '.tar', '.gz', '.rar'
    ]
    return ext in compatible_exts

def process_file(file_path, temp_dir, disposable_domains):
    ext = os.path.splitext(file_path)[1].lower()
    handlers = {
        '.txt': lambda f: read_text_file(f, disposable_domains),
        '.log': lambda f: read_text_file(f, disposable_domains),
        '.ini': lambda f: read_text_file(f, disposable_domains),
        '.inf': lambda f: read_text_file(f, disposable_domains),
        '.html': lambda f: read_text_file(f, disposable_domains),
        '.htm': lambda f: read_text_file(f, disposable_domains),
        '.asp': lambda f: read_text_file(f, disposable_domains),
        '.aspx': lambda f: read_text_file(f, disposable_domains),
        '.php': lambda f: read_text_file(f, disposable_domains),
        '.js': lambda f: read_text_file(f, disposable_domains),
        '.json': lambda f: read_text_file(f, disposable_domains),
        '.xml': lambda f: read_text_file(f, disposable_domains),
        '.yaml': lambda f: read_text_file(f, disposable_domains),
        '.yml': lambda f: read_text_file(f, disposable_domains),
        '.md': lambda f: read_text_file(f, disposable_domains),
        '.csv': lambda f: read_csv_file(f, disposable_domains),
        '.xls': lambda f: read_xls_file(f, disposable_domains),
        '.xlsx': lambda f: read_xlsx_file(f, disposable_domains),
        '.xlsm': lambda f: read_xlsx_file(f, disposable_domains),
        '.ods': lambda f: read_odt_file(f, disposable_domains),
        '.docx': lambda f: read_docx_file(f, disposable_domains),
        '.docm': lambda f: read_docx_file(f, disposable_domains),
        '.doc': lambda f: read_doc_file(f, disposable_domains),
        '.rtf': lambda f: read_rtf_file(f, disposable_domains),
        '.odt': lambda f: read_odt_file(f, disposable_domains),
        '.pptx': lambda f: read_pptx_file(f, disposable_domains),
        '.ppt': lambda f: read_ppt_file(f, disposable_domains),
        '.pdf': lambda f: read_pdf_file(f, disposable_domains),
        '.jpg': lambda f: read_image_file(f, disposable_domains),
        '.jpeg': lambda f: read_image_file(f, disposable_domains),
        '.png': lambda f: read_image_file(f, disposable_domains),
        '.bmp': lambda f: read_image_file(f, disposable_domains),
        '.tiff': lambda f: read_image_file(f, disposable_domains),
        '.gif': lambda f: read_image_file(f, disposable_domains),
        '.sqlite': lambda f: read_sqlite_file(f, disposable_domains),
        '.sqlite3': lambda f: read_sqlite_file(f, disposable_domains),
        '.db': lambda f: read_sqlite_file(f, disposable_domains),
        '.sql': lambda f: read_sql_file(f, disposable_domains),
        '.mdb': lambda f: read_mdb_file(f, disposable_domains),
        '.accdb': lambda f: read_mdb_file(f, disposable_domains),
        '.eml': lambda f: read_eml_file(f, disposable_domains),
        '.msg': lambda f: read_msg_file(f, disposable_domains),
        '.zip': lambda f: read_archive(f, temp_dir, disposable_domains),
        '.tar': lambda f: read_archive(f, temp_dir, disposable_domains),
        '.gz': lambda f: read_archive(f, temp_dir, disposable_domains),
        '.rar': lambda f: read_archive(f, temp_dir, disposable_domains),
    }
    if ext in handlers:
        return handlers[ext](file_path)
    kind = filetype.guess(file_path)
    if kind and kind.mime.startswith('image'):
        return read_image_file(file_path, disposable_domains)
    try:
        return read_text_file(file_path, disposable_domains)
    except Exception:
        return set(), ""

def scan_folder(folder, output_file, url_output_file, log_file, disposable_domains, blocked_domains):
    setup_logger(log_file)
    all_emails = set()
    all_urls = set()
    temp_dir = os.path.join(folder, "_temp_extract")
    os.makedirs(temp_dir, exist_ok=True)
    file_list = []
    for root, _, files in os.walk(folder):
        for file in files:
            if '_temp_extract' in root:
                continue
            path = os.path.join(root, file)
            file_list.append(path)
    total_files = len(file_list)
    compatible_files = [f for f in file_list if is_compatible_file(f)]
    total_compatible = len(compatible_files)

    print(f"Total files found: {total_files}")
    print(f"Compatible files for extraction: {total_compatible}")

    if total_compatible == 0:
        print("No compatible files found for extraction. Exiting.")
        logging.error("No compatible files found for extraction. Exiting.")
        return

    found_files = 0
    all_possible_emails = set()
    forbidden_filtered = set()
    disposable_filtered = set()
    exported_emails = set()
    all_possible_urls = set()
    blocked_urls = set()
    exported_urls = set()

    def classify_email(email):
        local, _, domain = email.lower().partition('@')
        if any(word in local or word in domain for word in FORBIDDEN_WORDS):
            return 'forbidden'
        domain_only = domain.split(':')[0].split('/')[0]
        if domain_only in disposable_domains:
            return 'disposable'
        return 'valid'

    with open(output_file, 'a', buffering=1) as f_out, \
         open(url_output_file, 'a', buffering=1) as f_url_out, \
         tqdm(total=total_compatible, desc="Extracting emails/urls", ncols=80) as pbar:
        for idx, path in enumerate(compatible_files, 1):
            print(f"Processing file {idx}/{total_compatible}: {path}")
            logging.info(f"Processing file {idx}/{total_compatible}: {path}")
            try:
                emails, text_content = process_file(path, temp_dir, disposable_domains)
                # Extract URLs from text_content
                urls = extract_urls_from_text(text_content, blocked_domains)
                for url in urls:
                    all_possible_urls.add(url)
                    parsed = urlparse(url)
                    domain = parsed.netloc.lower()
                    if domain.startswith('www.'):
                        domain = domain[4:]
                    if domain in blocked_domains:
                        blocked_urls.add(url)
                    elif url not in exported_urls:
                        f_url_out.write(url + '\n')
                        f_url_out.flush()
                        exported_urls.add(url)

                for email in emails:
                    all_possible_emails.add(email)
                    category = classify_email(email)
                    if category == 'forbidden':
                        forbidden_filtered.add(email)
                    elif category == 'disposable':
                        disposable_filtered.add(email)
                    elif email not in exported_emails:
                        f_out.write(email + '\n')
                        f_out.flush()
                        exported_emails.add(email)
                if emails or urls:
                    found_files += 1
                    print(f"  Found emails: {emails}")
                    print(f"  Found urls: {urls}")
                    logging.info(f"Found in {path}: {emails} {urls}")
            except Exception as e:
                print(f"  Error processing {path}: {e}")
                logging.error(f"Failed to process {path}: {e}")
            pbar.update(1)
    logging.info(f"Extraction complete. Unique emails found: {len(exported_emails)}")
    import shutil
    shutil.rmtree(temp_dir, ignore_errors=True)

    print("\n--- Extraction Summary ---")
    print(f"Total unique emails found (before filtering): {len(all_possible_emails)}")
    print(f"Removed due to forbidden words: {len(forbidden_filtered)}")
    print(f"Removed due to disposable domains: {len(disposable_filtered)}")
    print(f"Valid emails exported: {len(exported_emails)} (see {output_file})")
    print(f"Total unique urls found: {len(all_possible_urls)}")
    print(f"Removed due to blocked domains: {len(blocked_urls)}")
    print(f"Valid urls exported: {len(exported_urls)} (see {url_output_file})")
    print(f"Files with emails/urls found: {found_files} / {total_compatible}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Universal Email & URL Extractor (with Filtering and Summary)")
    parser.add_argument("folder", help="Folder to scan")
    parser.add_argument("-o", "--output", default="emails_found.txt", help="Output file for emails")
    parser.add_argument("-u", "--url_output", default="urls_found.txt", help="Output file for URLs")
    parser.add_argument("-l", "--log", default="extractor.log", help="Log file")
    parser.add_argument("-d", "--domains", default="disposable_domains.txt", help="Disposable domains file")
    parser.add_argument("-b", "--blocked_domains", default="blocked_domains.txt", help="Blocked URL domains text file (one domain per line)")
    args = parser.parse_args()
    disposable_domains = load_disposable_domains(args.domains)
    blocked_domains = load_blocked_domains(args.blocked_domains)
    scan_folder(args.folder, args.output, args.url_output, args.log, disposable_domains, blocked_domains)
